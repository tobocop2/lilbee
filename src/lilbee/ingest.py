"""Document sync engine — keeps documents/ dir in sync with LanceDB."""

import hashlib
import logging
import os
from collections.abc import Callable
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path
from typing import TypedDict

from rich.progress import Progress, SpinnerColumn, TextColumn

import lilbee.config as cfg
from lilbee import embedder, store
from lilbee.chunker import chunk_pages, chunk_text
from lilbee.code_chunker import CodeChunk, chunk_code, supported_extensions

log = logging.getLogger(__name__)


class ChunkRecord(TypedDict):
    """A single store-ready chunk record matching store._CHUNKS_SCHEMA."""

    source: str
    content_type: str
    page_start: int
    page_end: int
    line_start: int
    line_end: int
    chunk: str
    chunk_index: int
    vector: list[float]


def _build_records(chunks: list[str], source: str, content_type: str) -> list[ChunkRecord]:
    """Embed chunks and build store-ready records with zeroed page/line fields."""
    if not chunks:
        return []
    vectors = embedder.embed_batch(chunks)
    return [
        ChunkRecord(
            source=source,
            content_type=content_type,
            page_start=0,
            page_end=0,
            line_start=0,
            line_end=0,
            chunk=chunk,
            chunk_index=idx,
            vector=vec,
        )
        for idx, (chunk, vec) in enumerate(zip(chunks, vectors, strict=True))
    ]


# File extensions routed to the text chunker
_TEXT_EXTENSIONS = frozenset({".md", ".txt", ".html", ".rst"})

# File extensions routed to the code chunker
_CODE_EXTENSIONS = supported_extensions()

# Office document extensions
_OFFICE_EXTENSIONS: dict[str, str] = {".docx": "docx", ".xlsx": "xlsx", ".pptx": "pptx"}

# eBook extensions
_EBOOK_EXTENSIONS = frozenset({".epub"})

# Image extensions (OCR via Tesseract)
_IMAGE_EXTENSIONS = frozenset({".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp", ".webp"})

# Data file extensions
_DATA_EXTENSIONS = frozenset({".csv", ".tsv"})


def _file_hash(path: Path) -> str:
    """Compute SHA-256 hex digest of a file."""
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for block in iter(lambda: f.read(8192), b""):
            h.update(block)
    return h.hexdigest()


def _relative_name(path: Path) -> str:
    """Get path relative to documents dir as string."""
    return str(path.relative_to(cfg.DOCUMENTS_DIR))


def _discover_files() -> dict[str, Path]:
    """Scan documents/ recursively, return {relative_name: absolute_path}."""
    if not cfg.DOCUMENTS_DIR.exists():
        return {}
    files: dict[str, Path] = {}
    supported = (
        _TEXT_EXTENSIONS
        | _CODE_EXTENSIONS
        | frozenset({".pdf"})
        | frozenset(_OFFICE_EXTENSIONS)
        | _EBOOK_EXTENSIONS
        | _IMAGE_EXTENSIONS
        | _DATA_EXTENSIONS
    )
    for root, dirs, filenames in os.walk(cfg.DOCUMENTS_DIR, topdown=True):
        dirs[:] = [d for d in dirs if not cfg.is_ignored_dir(d)]
        for fname in filenames:
            if fname.startswith("."):
                continue
            path = Path(root) / fname
            if path.suffix.lower() in supported:
                files[_relative_name(path)] = path
    return files


def _classify_file(path: Path) -> str | None:
    """Classify file by extension. Returns content_type or None if unsupported."""
    ext = path.suffix.lower()
    if ext == ".pdf":
        return "pdf"
    if ext in _TEXT_EXTENSIONS:
        return "text"
    if ext in _CODE_EXTENSIONS:
        return "code"
    if ext in _OFFICE_EXTENSIONS:
        return _OFFICE_EXTENSIONS[ext]
    if ext in _EBOOK_EXTENSIONS:
        return "epub"
    if ext in _IMAGE_EXTENSIONS:
        return "image"
    if ext in _DATA_EXTENSIONS:
        return "data"
    return None


def _ingest_pdf(path: Path, source_name: str) -> list[ChunkRecord]:
    """Extract text from PDF, chunk, embed, and return store-ready records."""
    import pymupdf4llm

    md = pymupdf4llm.to_markdown(str(path), page_chunks=True)
    pages = [{"page": chunk["metadata"]["page"] + 1, "text": chunk["text"]} for chunk in md]

    page_chunks = chunk_pages(pages)
    if not page_chunks:
        return []

    texts = [pc.chunk for pc in page_chunks]
    vectors = embedder.embed_batch(texts)

    return [
        ChunkRecord(
            source=source_name,
            content_type="pdf",
            page_start=pc.page_start,
            page_end=pc.page_end,
            line_start=0,
            line_end=0,
            chunk=pc.chunk,
            chunk_index=pc.chunk_index,
            vector=vec,
        )
        for pc, vec in zip(page_chunks, vectors, strict=True)
    ]


def _ingest_text(path: Path, source_name: str) -> list[ChunkRecord]:
    """Read text file, chunk, embed, and return store-ready records."""
    text = path.read_text(encoding="utf-8", errors="replace")
    return _build_records(chunk_text(text), source_name, "text")


def _ingest_code(path: Path, source_name: str) -> list[ChunkRecord]:
    """Parse code with tree-sitter, chunk, embed, and return store-ready records."""
    code_chunks: list[CodeChunk] = chunk_code(path)
    if not code_chunks:
        return []

    texts = [cc.chunk for cc in code_chunks]
    vectors = embedder.embed_batch(texts)

    return [
        ChunkRecord(
            source=source_name,
            content_type="code",
            page_start=0,
            page_end=0,
            line_start=cc.line_start,
            line_end=cc.line_end,
            chunk=cc.chunk,
            chunk_index=cc.chunk_index,
            vector=vec,
        )
        for cc, vec in zip(code_chunks, vectors, strict=True)
    ]


def _ingest_docx(path: Path, source_name: str) -> list[ChunkRecord]:
    """Extract text from DOCX, chunk, embed, and return store-ready records."""
    from docx import Document

    doc = Document(str(path))
    parts: list[str] = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append("\t".join(cells))
    return _build_records(chunk_text("\n\n".join(parts)), source_name, "docx")


def _ingest_xlsx(path: Path, source_name: str) -> list[ChunkRecord]:
    """Extract text from XLSX, chunk, embed, and return store-ready records."""
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    parts: list[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        sheet_lines = [f"# Sheet: {sheet_name}"]
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                sheet_lines.append("\t".join(cells))
        if len(sheet_lines) > 1:
            parts.append("\n".join(sheet_lines))
    wb.close()
    return _build_records(chunk_text("\n\n".join(parts)), source_name, "xlsx")


def _ingest_pptx(path: Path, source_name: str) -> list[ChunkRecord]:
    """Extract text from PPTX, chunk, embed, and return store-ready records."""
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = [f"# Slide {slide_num}"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    slide_texts.append(text)
        if len(slide_texts) > 1:
            parts.append("\n".join(slide_texts))
    return _build_records(chunk_text("\n\n".join(parts)), source_name, "pptx")


def _ingest_epub(path: Path, source_name: str) -> list[ChunkRecord]:
    """Extract text from EPUB, chunk, embed, and return store-ready records."""
    import ebooklib
    from bs4 import BeautifulSoup
    from ebooklib import epub

    book = epub.read_epub(str(path))
    parts: list[str] = []
    for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
        html = item.get_content().decode("utf-8", errors="replace")
        text = BeautifulSoup(html, "html.parser").get_text(separator="\n").strip()
        if text:
            parts.append(text)
    return _build_records(chunk_text("\n\n".join(parts)), source_name, "epub")


def _ingest_image(path: Path, source_name: str) -> list[ChunkRecord]:
    """OCR an image, chunk, embed, and return store-ready records."""
    import pytesseract
    from PIL import Image

    text = pytesseract.image_to_string(Image.open(path)).strip()
    if not text:
        return []
    return _build_records(chunk_text(text), source_name, "image")


def _ingest_data(path: Path, source_name: str) -> list[ChunkRecord]:
    """Read CSV/TSV, chunk, embed, and return store-ready records."""
    import csv

    delimiter = "\t" if path.suffix.lower() == ".tsv" else ","
    with open(path, newline="", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f, delimiter=delimiter)
        lines = ["\t".join(row) for row in reader if any(cell.strip() for cell in row)]
    return _build_records(chunk_text("\n".join(lines)), source_name, "data")


_INGEST_DISPATCH: dict[str, Callable[[Path, str], list[ChunkRecord]]] = {
    "pdf": _ingest_pdf,
    "text": _ingest_text,
    "code": _ingest_code,
    "docx": _ingest_docx,
    "xlsx": _ingest_xlsx,
    "pptx": _ingest_pptx,
    "epub": _ingest_epub,
    "image": _ingest_image,
    "data": _ingest_data,
}


def _ingest_file(path: Path, source_name: str, content_type: str) -> int:
    """Ingest a single file. Returns chunk count."""
    ingest_fn = _INGEST_DISPATCH[content_type]
    records: list[dict] = ingest_fn(path, source_name)  # type: ignore[assignment]
    return store.add_chunks(records)


def sync(force_rebuild: bool = False, quiet: bool = False) -> dict:
    """Sync documents/ with the vector store.

    Returns summary dict with keys: added, updated, removed, unchanged, failed.
    When *quiet* is True, the Rich progress bar is suppressed (for JSON output).
    """
    if force_rebuild:
        store.drop_all()

    cfg.DOCUMENTS_DIR.mkdir(parents=True, exist_ok=True)

    disk_files = _discover_files()
    existing_sources = {s["filename"]: s["file_hash"] for s in store.get_sources()}

    added: list[str] = []
    updated: list[str] = []
    removed: list[str] = []
    unchanged = 0
    failed: list[str] = []

    # Find files to remove (in DB but not on disk)
    for name in existing_sources:
        if name not in disk_files:
            store.delete_by_source(name)
            store.delete_source(name)
            removed.append(name)

    # Process files on disk
    files_to_process: list[tuple[str, Path, str]] = []  # (name, path, content_type)

    for name, path in sorted(disk_files.items()):
        content_type = _classify_file(path)
        assert content_type is not None, f"Unsupported file slipped through discovery: {name}"

        current_hash = _file_hash(path)
        old_hash = existing_sources.get(name)

        if old_hash == current_hash:
            unchanged += 1
            continue

        if old_hash is not None:
            # Modified — remove old data
            store.delete_by_source(name)
            store.delete_source(name)
            files_to_process.append((name, path, content_type))
            updated.append(name)
        else:
            files_to_process.append((name, path, content_type))
            added.append(name)

    # Ingest files (with optional progress bar)
    if files_to_process:
        embedder.validate_model()
        _ingest_batch(files_to_process, added, updated, failed, quiet=quiet)

    return {
        "added": added,
        "updated": updated,
        "removed": removed,
        "unchanged": unchanged,
        "failed": failed,
    }


def _ingest_batch(
    files_to_process: list[tuple[str, Path, str]],
    added: list[str],
    updated: list[str],
    failed: list[str],
    *,
    quiet: bool = False,
) -> None:
    """Ingest a batch of files, optionally showing a Rich progress bar."""

    def _process_file(name: str, path: Path, content_type: str) -> tuple[str, int]:
        chunk_count = _ingest_file(path, name, content_type)
        return name, chunk_count

    workers = min(os.cpu_count() or 4, len(files_to_process))
    with ThreadPoolExecutor(max_workers=workers) as executor:
        futures = {
            executor.submit(_process_file, name, path, ct): (name, path)
            for name, path, ct in files_to_process
        }

        if quiet:
            _collect_results(futures, added, updated, failed)
        else:
            _collect_results_with_progress(futures, added, updated, failed)


def _collect_results(
    futures: dict,
    added: list[str],
    updated: list[str],
    failed: list[str],
) -> None:
    """Collect futures results without progress display."""
    for future in as_completed(futures):
        name, path = futures[future]
        try:
            _, chunk_count = future.result()
            store.upsert_source(name, _file_hash(path), chunk_count)
        except Exception:
            log.exception("Failed to ingest %s", name)
            if name in added:
                added.remove(name)
            if name in updated:
                updated.remove(name)
            failed.append(name)


def _collect_results_with_progress(
    futures: dict,
    added: list[str],
    updated: list[str],
    failed: list[str],
) -> None:
    """Collect futures results with Rich progress bar."""
    with Progress(
        SpinnerColumn(),
        TextColumn("{task.description}"),
        transient=True,
    ) as progress:
        task = progress.add_task("Ingesting documents...", total=len(futures))
        for future in as_completed(futures):
            name, path = futures[future]
            try:
                _, chunk_count = future.result()
                store.upsert_source(name, _file_hash(path), chunk_count)
            except Exception:
                log.exception("Failed to ingest %s", name)
                if name in added:
                    added.remove(name)
                if name in updated:
                    updated.remove(name)
                failed.append(name)
                progress.update(task, description=f"Failed {name}")
                progress.advance(task)
                continue
            progress.update(task, description=f"Ingested {name}")
            progress.advance(task)
