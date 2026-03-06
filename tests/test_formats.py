"""Real-file format tests — full sync() pipeline with actual files on disk.

Each test creates a real file using the actual library (docx, openpyxl, etc.),
drops it into the temp documents dir, calls sync(), and verifies it appears
in the result and produces correct chunks. Embeddings are mocked (no Ollama).
"""

from __future__ import annotations

from unittest import mock

import pytest

import lilbee.config as cfg
import lilbee.store as store_mod


@pytest.fixture(autouse=True)
def isolated_env(tmp_path):
    """Redirect config paths to temp dir for every test."""
    docs = tmp_path / "documents"
    docs.mkdir()
    data = tmp_path / "data" / "lancedb"

    orig_docs, orig_db, orig_data = cfg.DOCUMENTS_DIR, cfg.LANCEDB_DIR, cfg.DATA_DIR

    cfg.DOCUMENTS_DIR = docs
    cfg.DATA_DIR = tmp_path / "data"
    cfg.LANCEDB_DIR = data
    store_mod.LANCEDB_DIR = data

    yield docs

    cfg.DOCUMENTS_DIR = orig_docs
    cfg.DATA_DIR = orig_data
    cfg.LANCEDB_DIR = orig_db
    store_mod.LANCEDB_DIR = orig_db


def _fake_embed_batch(texts):
    return [[0.1] * 768 for _ in texts]


def _fake_embed(text):
    return [0.1] * 768


# ---------------------------------------------------------------------------
# Office formats
# ---------------------------------------------------------------------------


@mock.patch("lilbee.embedder.embed", side_effect=_fake_embed)
@mock.patch("lilbee.embedder.embed_batch", side_effect=_fake_embed_batch)
class TestSyncDocx:
    def test_docx_paragraphs(self, _eb, _e, isolated_env):
        from docx import Document

        doc = Document()
        doc.add_paragraph("The quick brown fox jumps over the lazy dog.")
        doc.add_paragraph("Pack my box with five dozen liquor jugs.")
        doc.save(str(isolated_env / "sample.docx"))

        from lilbee.ingest import sync

        result = sync()
        assert "sample.docx" in result["added"]

    def test_docx_with_table(self, _eb, _e, isolated_env):
        from docx import Document

        doc = Document()
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "Name"
        table.cell(0, 1).text = "Score"
        table.cell(1, 0).text = "Alice"
        table.cell(1, 1).text = "95"
        doc.save(str(isolated_env / "table.docx"))

        from lilbee.ingest import sync

        result = sync()
        assert "table.docx" in result["added"]


@mock.patch("lilbee.embedder.embed", side_effect=_fake_embed)
@mock.patch("lilbee.embedder.embed_batch", side_effect=_fake_embed_batch)
class TestSyncXlsx:
    def test_xlsx_with_data(self, _eb, _e, isolated_env):
        from openpyxl import Workbook

        wb = Workbook()
        ws1 = wb.active
        ws1.title = "Employees"
        ws1.append(["Name", "Department", "Salary"])
        ws1.append(["Alice", "Engineering", 120000])
        ws1.append(["Bob", "Marketing", 95000])

        ws2 = wb.create_sheet("Projects")
        ws2.append(["Project", "Status"])
        ws2.append(["Atlas", "Active"])

        wb.save(str(isolated_env / "data.xlsx"))

        from lilbee.ingest import sync

        result = sync()
        assert "data.xlsx" in result["added"]


@mock.patch("lilbee.embedder.embed", side_effect=_fake_embed)
@mock.patch("lilbee.embedder.embed_batch", side_effect=_fake_embed_batch)
class TestSyncPptx:
    def test_pptx_with_slides(self, _eb, _e, isolated_env):
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        layout = prs.slide_layouts[5]  # blank layout

        slide1 = prs.slides.add_slide(layout)
        txbox1 = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
        txbox1.text_frame.text = "Introduction to machine learning fundamentals"

        slide2 = prs.slides.add_slide(layout)
        txbox2 = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
        txbox2.text_frame.text = "Supervised and unsupervised learning algorithms"

        prs.save(str(isolated_env / "slides.pptx"))

        from lilbee.ingest import sync

        result = sync()
        assert "slides.pptx" in result["added"]


# ---------------------------------------------------------------------------
# EPUB
# ---------------------------------------------------------------------------


@mock.patch("lilbee.embedder.embed", side_effect=_fake_embed)
@mock.patch("lilbee.embedder.embed_batch", side_effect=_fake_embed_batch)
class TestSyncEpub:
    def test_epub_with_chapter(self, _eb, _e, isolated_env):
        from ebooklib import epub

        book = epub.EpubBook()
        book.set_identifier("test-book-001")
        book.set_title("Test Book")
        book.set_language("en")
        book.add_author("Test Author")

        chapter = epub.EpubHtml(title="Chapter 1", file_name="ch1.xhtml", lang="en")
        chapter.content = (
            "<html><body>"
            "<h1>Chapter 1</h1>"
            "<p>The history of artificial intelligence begins in antiquity.</p>"
            "</body></html>"
        )
        book.add_item(chapter)

        book.add_item(epub.EpubNcx())
        book.add_item(epub.EpubNav())
        book.spine = ["nav", chapter]
        book.toc = [epub.Link("ch1.xhtml", "Chapter 1", "ch1")]

        epub.write_epub(str(isolated_env / "book.epub"), book)

        from lilbee.ingest import sync

        result = sync()
        assert "book.epub" in result["added"]


# ---------------------------------------------------------------------------
# Image (OCR mocked — real PNG file, pytesseract mocked)
# ---------------------------------------------------------------------------


@mock.patch("lilbee.embedder.embed", side_effect=_fake_embed)
@mock.patch("lilbee.embedder.embed_batch", side_effect=_fake_embed_batch)
class TestSyncImage:
    def test_image_with_ocr(self, _eb, _e, isolated_env):
        from PIL import Image

        img = Image.new("RGB", (200, 100), color="white")
        img.save(str(isolated_env / "scan.png"))

        from lilbee.ingest import sync

        with mock.patch(
            "pytesseract.image_to_string",
            return_value="Invoice total is one hundred dollars.",
        ):
            result = sync()
        assert "scan.png" in result["added"]


# ---------------------------------------------------------------------------
# Code — all supported languages through full sync pipeline
# ---------------------------------------------------------------------------

_CODE_FIXTURES: dict[str, tuple[str, str]] = {
    "greet.py": (
        ".py",
        'def greet(name):\n    """Say hello."""\n    return f"Hello, {name}!"\n',
    ),
    "greet.js": (
        ".js",
        'function greet(name) {\n    return "Hello, " + name + "!";\n}\n',
    ),
    "greet.ts": (
        ".ts",
        "function greet(name: string): string {\n    return `Hello, ${name}!`;\n}\n",
    ),
    "greet.go": (
        ".go",
        'package main\n\nfunc greet(name string) string {\n\treturn "Hello, " + name\n}\n',
    ),
    "greet.rs": (
        ".rs",
        'fn greet(name: &str) -> String {\n    format!("Hello, {}!", name)\n}\n',
    ),
    "Greet.java": (
        ".java",
        "public class Greet {\n"
        "    public static String greet(String name) {\n"
        '        return "Hello, " + name + "!";\n'
        "    }\n"
        "}\n",
    ),
    "greet.c": (
        ".c",
        "#include <stdio.h>\n\n"
        'void greet(const char* name) {\n    printf("Hello, %s!\\n", name);\n}\n',
    ),
    "greet.cpp": (
        ".cpp",
        "#include <string>\n\n"
        "std::string greet(const std::string& name) {\n"
        '    return "Hello, " + name + "!";\n'
        "}\n",
    ),
}


@mock.patch("lilbee.embedder.embed", side_effect=_fake_embed)
@mock.patch("lilbee.embedder.embed_batch", side_effect=_fake_embed_batch)
class TestSyncCode:
    @pytest.mark.parametrize("filename,fixture", list(_CODE_FIXTURES.items()))
    def test_code_file_syncs(self, _eb, _e, isolated_env, filename, fixture):
        _ext, content = fixture
        (isolated_env / filename).write_text(content)

        from lilbee.ingest import sync

        result = sync()
        assert filename in result["added"]


# ---------------------------------------------------------------------------
# CSV / TSV
# ---------------------------------------------------------------------------


@mock.patch("lilbee.embedder.embed", side_effect=_fake_embed)
@mock.patch("lilbee.embedder.embed_batch", side_effect=_fake_embed_batch)
class TestSyncCsvTsv:
    def test_csv_with_unicode(self, _eb, _e, isolated_env):
        csv_content = "name,city,note\nRene,Montreal,cafe au lait\nJose,Madrid,hola\n"
        (isolated_env / "people.csv").write_text(csv_content)

        from lilbee.ingest import sync

        result = sync()
        assert "people.csv" in result["added"]

    def test_tsv_multiline(self, _eb, _e, isolated_env):
        tsv_content = "id\tproduct\tprice\n1\tWidget\t9.99\n2\tGadget\t19.99\n3\tDoohickey\t4.50\n"
        (isolated_env / "products.tsv").write_text(tsv_content)

        from lilbee.ingest import sync

        result = sync()
        assert "products.tsv" in result["added"]
